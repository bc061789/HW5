import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import os

st.set_page_config(page_title="AI PPT Auto Formatter", layout="centered")

st.title("📊 AI PPT 自動版型調整系統")
st.write("上傳你的 PPT，系統會自動重新排版與統一風格")

uploaded_file = st.file_uploader("請上傳 PPT 檔案", type=["pptx"])

def auto_format_ppt(input_path, output_path):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)

                paragraph.alignment = PP_ALIGN.LEFT

    prs.save(output_path)

if uploaded_file:
    st.success("✅ 上傳成功，開始自動調整版型中...")

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_input:
        temp_input.write(uploaded_file.read())
        input_path = temp_input.name

    output_path = input_path.replace(".pptx", "_new.pptx")

    auto_format_ppt(input_path, output_path)

    with open(output_path, "rb") as f:
        st.download_button(
            label="📥 下載新版 PPT",
            data=f,
            file_name="formatted_ppt.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    os.remove(input_path)
    os.remove(output_path)
